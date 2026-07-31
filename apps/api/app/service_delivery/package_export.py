import csv
import hashlib
import io
import json
import re
import zipfile
from datetime import datetime
from typing import Any

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


MIME_TYPES = {
    "md": "text/markdown",
    "json": "application/json",
    "csv": "text/csv",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "zip": "application/zip",
}


def _mime_type(path: str) -> str:
    extension = path.rsplit(".", 1)[-1].casefold() if "." in path else ""
    return MIME_TYPES.get(extension, {
        "py": "text/x-python",
        "ts": "text/typescript",
        "tsx": "text/typescript",
        "js": "text/javascript",
        "jsx": "text/javascript",
        "css": "text/css",
        "html": "text/html",
        "toml": "application/toml",
        "yaml": "application/yaml",
        "yml": "application/yaml",
        "txt": "text/plain",
    }.get(extension, "application/octet-stream"))


def _stable_zip(payload: bytes) -> bytes:
    source = zipfile.ZipFile(io.BytesIO(payload), "r")
    output = io.BytesIO()
    with source, zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as target:
        for name in sorted(source.namelist()):
            info = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            data = source.read(name)
            if name == "docProps/core.xml":
                # openpyxl replaces `modified` during save even when callers
                # set a fixed property. Normalize OOXML metadata after save.
                data = re.sub(
                    rb"(<dcterms:modified\b[^>]*>).*?(</dcterms:modified>)",
                    rb"\g<1>1980-01-01T00:00:00Z\g<2>",
                    data,
                )
            target.writestr(info, data)
    return output.getvalue()


def _safe_name(value: str) -> str:
    return re.sub(r"[^a-zA-Z0-9._-]+", "-", value).strip("-")[:100] or "deliverable"


def _markdown(content: dict[str, Any], fallback_title: str) -> str:
    markdown = str(content.get("content_markdown") or "").strip()
    if markdown:
        return markdown + "\n"
    title = str(content.get("title") or fallback_title)
    body = json.dumps(content, ensure_ascii=False, indent=2, default=str)
    return f"# {title}\n\n```json\n{body}\n```\n"


def _docx(title: str, markdown: str) -> bytes:
    document = Document()
    document.core_properties.created = datetime(1980, 1, 1)
    document.core_properties.modified = datetime(1980, 1, 1)
    document.add_heading(title, level=0)
    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped.startswith("### "):
            document.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            document.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            document.add_heading(stripped[2:], level=1)
        elif stripped.startswith("- "):
            document.add_paragraph(stripped[2:], style="List Bullet")
        elif stripped:
            document.add_paragraph(stripped)
    buffer = io.BytesIO()
    document.save(buffer)
    return _stable_zip(buffer.getvalue())


def _pptx(title: str, markdown: str) -> bytes:
    presentation = Presentation()
    presentation.core_properties.created = datetime(1980, 1, 1)
    presentation.core_properties.modified = datetime(1980, 1, 1)
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = title
    cover.placeholders[1].text = "Pacote editável gerado pela Agentic Software Factory"
    chunks = [line.strip("# -") for line in markdown.splitlines() if line.strip()][:30]
    for offset in range(0, len(chunks), 6):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = chunks[offset] or title
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.8))
        frame = textbox.text_frame
        frame.text = "\n".join(chunks[offset + 1:offset + 6])
    buffer = io.BytesIO()
    presentation.save(buffer)
    return _stable_zip(buffer.getvalue())


def _xlsx(title: str, content: dict[str, Any], evidence_refs: list[str]) -> bytes:
    workbook = Workbook()
    workbook.properties.created = datetime(1980, 1, 1)
    workbook.properties.modified = datetime(1980, 1, 1)
    sheet = workbook.active
    sheet.title = "Deliverable"
    sheet.append(["Campo", "Valor"])
    sheet.append(["Título", title])
    for key, value in content.items():
        sheet.append([str(key), json.dumps(value, ensure_ascii=False, default=str) if not isinstance(value, str) else value])
    evidence = workbook.create_sheet("Evidências")
    evidence.append(["Referência"])
    for ref in evidence_refs:
        evidence.append([ref])
    # openpyxl refreshes `modified` while saving. Reset it only after all
    # worksheet mutations so identical inputs produce byte-identical packages.
    workbook.properties.modified = datetime(1980, 1, 1)
    buffer = io.BytesIO()
    workbook.save(buffer)
    return _stable_zip(buffer.getvalue())


def _csv(content: dict[str, Any]) -> bytes:
    buffer = io.StringIO()
    writer = csv.writer(buffer)
    writer.writerow(["field", "value"])
    for key, value in content.items():
        writer.writerow([key, json.dumps(value, ensure_ascii=False, default=str) if not isinstance(value, str) else value])
    return buffer.getvalue().encode("utf-8")


def build_deliverable_package(
    *,
    deliverable: dict[str, Any],
    revision: dict[str, Any],
    formats: list[str],
    technical_files: dict[str, bytes] | None = None,
    technical_evidence: dict[str, Any] | None = None,
) -> tuple[str, bytes, dict[str, Any]]:
    title = str(deliverable["title"])
    stem = _safe_name(title)
    content = dict(revision.get("content_json") or {})
    evidence_refs = list(revision.get("evidence_refs_json") or [])
    markdown = _markdown(content, title)
    files: dict[str, bytes] = {
        f"sources/{stem}.md": markdown.encode("utf-8"),
        f"sources/{stem}.json": json.dumps(content, ensure_ascii=False, indent=2, default=str).encode("utf-8"),
        "evidence/evidence.json": json.dumps(evidence_refs, ensure_ascii=False, indent=2).encode("utf-8"),
    }
    for output_format in formats:
        if output_format == "docx":
            files[f"editable/{stem}.docx"] = _docx(title, markdown)
        elif output_format == "pptx":
            files[f"editable/{stem}.pptx"] = _pptx(title, markdown)
        elif output_format == "xlsx":
            files[f"editable/{stem}.xlsx"] = _xlsx(title, content, evidence_refs)
        elif output_format == "csv":
            files[f"editable/{stem}.csv"] = _csv(content)
    for path, data in sorted((technical_files or {}).items()):
        files[f"technical/source/{path}"] = data
    if technical_evidence:
        files["technical/evidence.json"] = json.dumps(
            technical_evidence,
            ensure_ascii=False,
            indent=2,
            default=str,
        ).encode("utf-8")
    entries = []
    for path, data in sorted(files.items()):
        entries.append(
            {
                "path": path,
                "sha256": hashlib.sha256(data).hexdigest(),
                "mime_type": _mime_type(path),
                "size_bytes": len(data),
                "origin": f"service_deliverable:{deliverable['id']}",
                "revision": revision["revision"],
            }
        )
    manifest = {
        "schema_version": "service-delivery-package/2.1" if technical_files else "service-delivery-package/2.0",
        "deliverable_id": deliverable["id"],
        "engagement_id": deliverable["engagement_id"],
        "revision": revision["revision"],
        "files": entries,
    }
    files["manifest.json"] = json.dumps(manifest, ensure_ascii=False, indent=2).encode("utf-8")
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for path, data in sorted(files.items()):
            info = zipfile.ZipInfo(path, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            archive.writestr(info, data)
    payload = output.getvalue()
    manifest["package_sha256"] = hashlib.sha256(payload).hexdigest()
    manifest["package_size_bytes"] = len(payload)
    return f"{stem}-r{revision['revision']}.zip", payload, manifest


def build_engagement_package(
    *,
    engagement: dict[str, Any],
    offering_version: str,
    deliverable_packages: list[tuple[str, bytes, dict[str, Any]]],
    acceptance_checks: list[dict[str, Any]],
) -> tuple[str, bytes, dict[str, Any]]:
    stem = _safe_name(str(engagement["name"]))
    files: dict[str, bytes] = {}
    package_refs: list[dict[str, Any]] = []
    for filename, payload, manifest in sorted(deliverable_packages, key=lambda item: item[0]):
        path = f"deliverables/{filename}"
        files[path] = payload
        package_refs.append({
            "path": path,
            "deliverable_id": manifest["deliverable_id"],
            "sha256": hashlib.sha256(payload).hexdigest(),
            "size_bytes": len(payload),
        })
    evidence = {
        "engagement": engagement,
        "offering_version": offering_version,
        "acceptance_checks": acceptance_checks,
        "deliverable_packages": package_refs,
    }
    files["evidence/engagement.json"] = json.dumps(
        evidence,
        ensure_ascii=False,
        indent=2,
        default=str,
    ).encode("utf-8")
    entries = [
        {
            "path": path,
            "sha256": hashlib.sha256(data).hexdigest(),
            "mime_type": _mime_type(path),
            "size_bytes": len(data),
            "origin": f"engagement:{engagement['id']}",
            "version": offering_version,
        }
        for path, data in sorted(files.items())
    ]
    manifest = {
        "schema_version": "engagement-delivery-package/2.1",
        "engagement_id": engagement["id"],
        "offering_version": offering_version,
        "files": entries,
    }
    files["manifest.json"] = json.dumps(manifest, ensure_ascii=False, indent=2).encode("utf-8")
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for path, data in sorted(files.items()):
            info = zipfile.ZipInfo(path, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o600 << 16
            archive.writestr(info, data)
    payload = output.getvalue()
    manifest["package_sha256"] = hashlib.sha256(payload).hexdigest()
    manifest["package_size_bytes"] = len(payload)
    return f"{stem}-portfolio-{offering_version}.zip", payload, manifest
